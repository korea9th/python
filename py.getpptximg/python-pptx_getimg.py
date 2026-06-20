import os
from pptx import Presentation

def extract_images_from_pptx(fullname, filename):
    # 출력 폴더가 없으면 생성
    if not os.path.exists(filename):
        os.makedirs(filename)    

    title = filename.split(' ')

    # 프레젠테이션 불러오기
    prs = Presentation(fullname)
    image_count = 0
   
    # 각 슬라이드 순회
    for slide_idx, slide in enumerate(prs.slides):
        image_count = 0
        for shape in slide.shapes:
            # 도형이 그림(이미지) 타입인지 확인
            if shape.shape_type == 13: # 13은 Picture 타입
                try:
                    # 이미지 객체 가져오기
                    image = shape.image
                    image_bytes = image.blob
                    image_ext = image.ext # 확장자 (예: 'jpeg', 'png')

                    # 저장할 파일 경로 설정
                    image_filename = f"{title[0]} Slide{slide_idx + 1:03d}_Img{image_count + 1:03d}.{image_ext}"
                    image_path = os.path.join(filename, image_filename)

                    # 파일 쓰기
                    with open(image_path, "wb") as f:
                        f.write(image_bytes)
                    
#                    print(f"저장 완료: {image_filename}")
                    image_count += 1
                    
                except Exception as e:
                    print(f"이미지 추출 실패: {e}")

    print(f"총 {image_count}개의 이미지가 '{filename}' 폴더에 저장되었습니다.")

# 사용 예시
#extract_images_from_pptx("T_INS.1 LogCops V6.0 설치 및 시동 시험.pptx", "extracted_images")


#send it
def search(dirname):
    global file_count, line_count, error_count
    filenames = os.listdir(dirname)
    for filename in filenames:
        fullname = os.path.join(dirname, filename)
        if os.path.isdir(fullname):
#            print("Dir : " + fullname)
            search(fullname)
        else:
            extract_images_from_pptx(fullname, filename)
#            getImg(fullname, filename)



search("C:\\02.dev\\getpptximg\\pptx")
